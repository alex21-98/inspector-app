import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from adjustText import adjust_text
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import io
import datetime

# ==========================================
# CONFIGURACIÓN DE LA PÁGINA
# ==========================================
st.set_page_config(page_title="Reporte MP-Calidad", layout="wide", page_icon="📊")

st.title("📊🖥️ Generador de Reportes de MP-Calidad Gerencial")
st.markdown("Complejo Agroindustrial Beta - Sistema de Visualización de Defectos")

# ==========================================
# 1. CARGA DE DATOS Y PREPARACIÓN
# ==========================================
MAPEO_COLUMNAS = {
    'Fecha': ['fecha', 'fec', 'date'], 'Mes': ['mes', 'month'],
    'Año': ['año', 'anio', 'year'], 'Semana': ['semana','SEMANA', 'sem', 'week', 'no sem'],
    'Fundo': ['Fundos', 'campo', 'fundo'], 'Lote': ['lote', 'lot', 'id lote'],
    'Variedad': ['variedad', 'var']
}

@st.cache_data
def cargar_y_limpiar_datos(file, sheet, freq):
    df_raw = pd.read_excel(file, sheet_name=sheet)
    df_proc = df_raw.copy()
    nuevos_nombres = {}
    
    for col_excel in df_proc.columns:
        col_limpia = str(col_excel).strip().lower()
        for nombre_estandar, sinonimos in MAPEO_COLUMNAS.items():
            if col_limpia in [s.lower() for s in sinonimos] or col_limpia == nombre_estandar.lower():
                nuevos_nombres[col_excel] = nombre_estandar
                break

    df_proc = df_proc.rename(columns=nuevos_nombres)
    df_proc = df_proc.loc[:, ~df_proc.columns.duplicated()]
    columnas_defectos = [col for col in df_proc.columns if '%' in str(col)]

    if 'Fecha' in df_proc.columns:
        df_proc['Fecha'] = pd.to_datetime(df_proc['Fecha'], errors='coerce', dayfirst=True)
        df_proc = df_proc.dropna(subset=['Fecha'])

        # Lógica de Periodos
        if freq == 'D':
            df_proc['Periodo'] = df_proc['Fecha'].dt.strftime('%d/%m/%Y')
            df_proc['Orden_Periodo'] = df_proc['Fecha'].dt.normalize()
        elif freq == 'W':
            df_proc['Orden_Periodo'] = df_proc['Fecha'] - pd.to_timedelta(df_proc['Fecha'].dt.dayofweek, unit='d')
            if 'Semana' in df_proc.columns:
                df_proc['Periodo'] = 'Sem ' + df_proc['Semana'].fillna(0).astype(int).astype(str)
            else:
                df_proc['Periodo'] = df_proc['Orden_Periodo'].dt.strftime('Sem %V')
        elif freq == 'M':
            df_proc['Periodo'] = df_proc['Fecha'].dt.strftime('%m/%Y')
            df_proc['Orden_Periodo'] = df_proc['Fecha'] - pd.to_timedelta(df_proc['Fecha'].dt.day - 1, unit='d')
        elif freq == 'Y':
            df_proc['Periodo'] = df_proc['Fecha'].dt.strftime('%Y')
            df_proc['Orden_Periodo'] = pd.to_datetime(df_proc['Fecha'].dt.year, format='%Y')

    columnas_requeridas = ['Periodo', 'Orden_Periodo', 'Fundo', 'Lote', 'Variedad']

    if all(col in df_proc.columns for col in columnas_requeridas):
        df_proc['Variedad'] = df_proc['Variedad'].astype(str)
        df_proc['Fundo'] = df_proc['Fundo'].astype(str)
        df_proc['Lote'] = df_proc['Lote'].astype(str)

        columnas_agrupar = columnas_requeridas + columnas_defectos
        df_agrupado = df_proc[columnas_agrupar].groupby(columnas_requeridas).mean().reset_index()
        df_agrupado = df_agrupado.sort_values(by=columnas_requeridas, ascending=[True, True, True, True, True])
        return df_agrupado, columnas_defectos, df_proc
    else:
        st.error("❌ Faltan columnas críticas.")
        return pd.DataFrame(), [], pd.DataFrame()

# --- PANEL LATERAL ---
with st.sidebar:
    st.header("📁 1. Carga de Datos")
    uploaded_file = st.file_uploader("Sube tu archivo Excel (.xlsx)", type=["xlsx"])
    if uploaded_file is not None:
        excel_obj = pd.ExcelFile(uploaded_file)
        hoja_final = st.selectbox("👉 Selecciona la hoja:", excel_obj.sheet_names)
        st.header("⚙️ 2. Periodicidad")
        opcion_freq = st.selectbox("Agrupar reporte por:", ["Semanal", "Diario", "Mensual", "Anual"])
        mapa_freq = {"Diario": 'D', "Semanal": 'W', "Mensual": 'M', "Anual": 'Y'}
        freq_elegida = mapa_freq[opcion_freq]

# --- LÓGICA PRINCIPAL ---
if uploaded_file is not None:
    df_final, lista_defectos, df_raw_proc = cargar_y_limpiar_datos(uploaded_file, hoja_final, freq_elegida)
    
    if not df_final.empty:
        st.divider()
        col_f1, col_f2 = st.columns(2)
        fecha_min = df_raw_proc['Fecha'].min().date()
        fecha_max = df_raw_proc['Fecha'].max().date()
        
        fecha_ini = col_f1.date_input("Fecha Inicio", value=fecha_min, min_value=fecha_min, max_value=fecha_max)
        fecha_fin = col_f2.date_input("Fecha Fin", value=fecha_max, min_value=fecha_min, max_value=fecha_max)
        
        mask = (df_final['Orden_Periodo'].dt.date >= fecha_ini) & (df_final['Orden_Periodo'].dt.date <= fecha_fin)
        df_plot = df_final[mask].copy()
        
        col1, col2, col3 = st.columns(3)
        fundos_sel = col1.multiselect("Fundos", df_plot['Fundo'].unique())
        if fundos_sel: df_plot = df_plot[df_plot['Fundo'].isin(fundos_sel)]
        variedades_sel = col2.multiselect("Variedades", df_plot['Variedad'].unique())
        if variedades_sel: df_plot = df_plot[df_plot['Variedad'].isin(variedades_sel)]
        df_plot['Etiqueta_Lote'] = df_plot['Fundo'] + " - " + df_plot['Lote']
        lotes_sel = col3.multiselect("Lotes", df_plot['Etiqueta_Lote'].unique())
        if lotes_sel: df_plot = df_plot[df_plot['Etiqueta_Lote'].isin(lotes_sel)]
        
        defectos_sel = st.multiselect("Defectos a graficar", lista_defectos)

        # ==========================================
        # FILTRO PARA NIVEL DE DETALLE DEL GRÁFICO
        # ==========================================
        st.divider()
        st.subheader("📈 3. Configuración Visual de Gráficos")
        col_visual1, col_visual2 = st.columns(2)
        
        with col_visual1:
            nivel_agrupacion = st.radio(
                "Selecciona cómo deseas agrupar las líneas en el gráfico:",
                ["Por Lote", "Por Fundo"],
                horizontal=True
            )
            col_agrupacion = 'Etiqueta_Lote' if nivel_agrupacion == "Por Lote" else 'Fundo'
            
        with col_visual2:
            titulo_personalizado = st.text_input(
                "Título base del gráfico:", 
                value="Evaluación de PT",
                help="Cambia este texto para modificar el inicio del título en gráficos y diapositivas"
            )

        # ==========================================
        # TOLERANCIAS DINÁMICAS
        # ==========================================
        tolerancias_defectos = {}
        if defectos_sel:
            st.divider()
            st.subheader("📏 4. Configuración de Tolerancias (%)")
            st.write("Establece el límite permitido para cada defecto. Déjalo en 0 si no deseas línea de tolerancia.")
            
            cols = st.columns(min(len(defectos_sel), 4))
            for i, defecto in enumerate(defectos_sel):
                with cols[i % 4]:
                    tol_val = st.number_input(f"{defecto}", min_value=0.0, max_value=100.0, value=0.0, step=0.5, key=f"tol_{defecto}")
                    if tol_val > 0:
                        tolerancias_defectos[defecto] = tol_val

        # ==========================================
        # GENERACIÓN Y VISUALIZACIÓN
        # ==========================================
        st.divider()
        if st.button("🚀 Generar Reporte Gerencial", type="primary", use_container_width=True):
            if not fundos_sel or not variedades_sel or not lotes_sel or not defectos_sel:
                st.warning("⚠️ Debes seleccionar al menos un Fundo, Variedad, Lote y Defecto para generar el reporte.")
            else:
                prs = Presentation()
                diapositivas_creadas = 0
                
                color_texto_principal = '#45605A'
                color_borde_grafico = '#B0BEC5'
                
                colores_fuertes = [
                    '#1976D2', '#388E3C', '#FBC02D', '#8E24AA', '#F57C00', '#0097A7', '#689F38', '#C2185B', '#111111', '#455A64',
                    '#3F51B5', '#00796B', '#AFB42B', '#512DA8', '#0288D1', '#F50057', '#C0CA33', '#8D6E63', '#26A69A', '#00ACC1',
                    '#2979FF', '#00C853', '#FFAB00', '#D500F9', '#FF6D00', '#00B8D4', '#AEEA00', '#C51162', '#607D8B', '#FF4081',
                    '#795548', '#311B92', '#004D40', '#827717', '#3E2723', '#01579B', '#1B5E20', '#E65100', '#4A148C', '#263238'
                ]

                graficos_expander = st.expander("Ver Previsualización de Gráficos", expanded=True)

                for defecto in defectos_sel:
                    for var in variedades_sel:
                        data_var_raw = df_plot[df_plot['Variedad'] == var]
                        if data_var_raw[defecto].isnull().all() or data_var_raw.empty: continue

                        cols_groupby = [col_agrupacion, 'Periodo', 'Orden_Periodo']
                        if 'Fundo' not in cols_groupby:
                            cols_groupby.append('Fundo')
                        data_var = data_var_raw.groupby(cols_groupby, dropna=False)[defecto].mean().reset_index()

                        fig, ax = plt.subplots(figsize=(14, 8), dpi=150)
                        
                        periodos_ordenados = data_var.sort_values('Orden_Periodo')['Periodo'].unique()
                        
                        entidades_presentes = data_var[col_agrupacion].unique()
                        textos_a_ajustar = []

                        max_val_test = data_var[defecto].max()
                        es_escala_decimal = max_val_test < 1.0 if pd.notna(max_val_test) else False

                        # =========================================================
                        # NUEVO: Lógica dinámica de tamaño de texto según cantidad
                        # =========================================================
                        num_lineas = len(entidades_presentes)
                        if num_lineas <= 3:
                            fs_dinamico = 14
                            pad_box = 0.3
                        elif num_lineas <= 6:
                            fs_dinamico = 12
                            pad_box = 0.25
                        elif num_lineas <= 10:
                            fs_dinamico = 10
                            pad_box = 0.2
                        else:
                            fs_dinamico = 9
                            pad_box = 0.15

                        for i, entidad in enumerate(entidades_presentes):
                            data_entidad = data_var[data_var[col_agrupacion] == entidad]
                            color_asignado = colores_fuertes[i % len(colores_fuertes)]

                            valores_entidad_alineados = []
                            for per in periodos_ordenados:
                                valor_per = data_entidad.loc[data_entidad['Periodo'] == per, defecto]
                                if not valor_per.empty and pd.notna(valor_per.iloc[0]):
                                    valores_entidad_alineados.append(valor_per.iloc[0])
                                else:
                                    valores_entidad_alineados.append(np.nan)

                            ax.plot(periodos_ordenados, valores_entidad_alineados, marker='o', label=entidad, 
                                    linewidth=4, color=color_asignado, markersize=8, markeredgecolor='white', 
                                    markeredgewidth=1.5, zorder=5)

                            for x_val, p in zip(periodos_ordenados, valores_entidad_alineados):
                                if pd.notna(p):
                                    val_etq = p * 100 if es_escala_decimal else p
                                    # CAMBIO: Usamos fs_dinamico y pad_box
                                    t = ax.text(x_val, p, f'{val_etq:.1f}%', fontsize=fs_dinamico, fontweight='bold', color='white', 
                                                ha='center', va='center',
                                                zorder=10, bbox=dict(facecolor=color_asignado, alpha=0.9, edgecolor='white', boxstyle=f'square,pad={pad_box}'))
                                    textos_a_ajustar.append(t)

                        if textos_a_ajustar:
                            # CAMBIO IMPORTANTE: Ajuste de colisiones inteligente
                            # Permitimos movimiento XY, pero la fuerza horizontal (force_text X=1.5) 
                            # es mucho mayor que la vertical (force_text Y=0.2).
                            # Esto evita que un valor de 5% salte por debajo del 2%, mandándolo a un costado con una línea que lo señala claramente.
                            adjust_text(
                                textos_a_ajustar, 
                                ax=ax, 
                                only_move={'points': 'xy', 'text': 'xy'}, 
                                expand_points=(1.2, 1.2), 
                                expand_text=(1.1, 1.1),
                                force_text=(1.5, 0.2),    # Fuerte empuje horizontal, muy leve vertical
                                force_points=(1.5, 0.2),  # Igual para la repulsión contra los puntos
                                arrowprops=dict(arrowstyle='-', color='#78909C', lw=1.2, alpha=0.8, zorder=2), 
                                max_move=40               # Limita qué tan lejos pueden volar las etiquetas
                            )

                        if defecto in tolerancias_defectos:
                            valor_tol = tolerancias_defectos[defecto]
                            valor_tol_grafico = valor_tol / 100 if es_escala_decimal else valor_tol
                            ax.axhline(y=valor_tol_grafico, color='#D32F2F', linestyle='--', linewidth=2.5, label=f'Límite ({valor_tol}%)', zorder=3)
                            x_pos_final = len(periodos_ordenados) - 1 if len(periodos_ordenados) > 0 else 0
                            ax.text(0, valor_tol_grafico, f' T: {valor_tol}% ', color='white', fontsize=12, fontweight='bold',
                                    ha='left', va='bottom', zorder=15, bbox=dict(facecolor='#D32F2F', edgecolor='white', alpha=0.9, boxstyle='round,pad=0.3'))

                        columnas_leyenda = min(len(entidades_presentes) + 1, 5)
                        
                        total_items = len(entidades_presentes) + 1
                        filas_leyenda = (total_items + columnas_leyenda - 1) // columnas_leyenda

                        pad_dinamico = 10 + (filas_leyenda * 15)

                        ax.legend(loc='lower center', bbox_to_anchor=(0.5, 1.02), ncol=columnas_leyenda, frameon=False, fontsize=11)
                        
                        texto_fundos = " y ".join(data_var['Fundo'].unique())
                        
                        ax.set_title(f"{titulo_personalizado}: {defecto} - {texto_fundos}\n".upper(), fontsize=18, fontweight='bold', color=color_texto_principal, pad=pad_dinamico)
                        
                        ax.set_xlabel(f"\nVariedad: {str(var).upper()}", fontsize=14, fontweight='bold', color=color_texto_principal)
                        ax.set_xticklabels(periodos_ordenados, rotation=45, ha='right', fontsize=12)
                        ax.get_yaxis().set_visible(False)
                        ax.margins(y=0.25)
                        
                        for spine in ax.spines.values():
                            spine.set_visible(True)
                            spine.set_color(color_borde_grafico)
                            spine.set_linewidth(2)

                        plt.figtext(0.99, 0.01, 'Complejo Agroindustrial Beta', horizontalalignment='right', fontsize=10, color='gray', style='italic')
                        plt.tight_layout()

                        graficos_expander.pyplot(fig)

                        image_stream = io.BytesIO()
                        fig.savefig(image_stream, format='png', dpi=300, bbox_inches='tight', transparent=False)
                        plt.close(fig)

                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        title_shape = slide.shapes.title
                        title_shape.text = f"{titulo_personalizado}: {defecto}"
                        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(69, 96, 90)
                        title_shape.text_frame.paragraphs[0].font.bold = True
                        
                        image_stream.seek(0)
                        slide.shapes.add_picture(image_stream, Inches(0.25), Inches(1.8), width=Inches(9.5))
                        diapositivas_creadas += 1

                if diapositivas_creadas > 0:
                    pptx_stream = io.BytesIO()
                    prs.save(pptx_stream)
                    pptx_stream.seek(0)
                    st.success(f"✅ Reporte generado exitosamente con {diapositivas_creadas} diapositivas.")
                    st.download_button(
                        label="📥 Descargar Presentación (PPTX)",
                        data=pptx_stream,
                        file_name=f"Reporte_Calidad_{datetime.datetime.now().strftime('%Y%m%d')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        type="primary"
                    )
else:
    st.info("👋 Sube un archivo Excel en el menú lateral para comenzar a configurar el reporte.")
